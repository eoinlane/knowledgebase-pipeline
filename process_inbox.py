#!/usr/bin/env python3
"""
process_inbox.py — Process files dropped into ~/inbox/

Supported formats:
  .pdf   — text extracted with pypdf
  .docx  — text extracted with python-docx
  .eml   — email + embedded attachments (PDF/docx extracted inline)
  .txt   — read directly
  .md    — read directly
  .png/.jpg/.jpeg/.gif/.webp — sent to Claude vision via LiteLLM

For each file:
  1. Extract text (or use vision for images)
  2. Classify with Claude via LiteLLM proxy (category/topic/summary/key_people)
  3. Write KB markdown to ~/knowledge_base/documents/
  4. Move original to ~/inbox/done/
  5. Trigger incremental upload to Open WebUI

Run: python3 ~/knowledgebase-pipeline/process_inbox.py
Also triggered by launchd WatchPaths on ~/inbox/
"""

import base64
import email
import html
import io
import json
import re
import shutil
import subprocess
from datetime import datetime
from email import policy as email_policy
from pathlib import Path

import requests

LITELLM_URL   = "http://100.121.184.27:4000/chat/completions"
MODEL         = "claude-haiku-4-5"

INBOX_DIR     = Path.home() / "inbox"
DONE_DIR      = INBOX_DIR / "done"
DOCS_DIR      = Path.home() / "knowledge_base" / "documents"
UPLOAD_SCRIPT = Path.home() / "upload_knowledge_base_incremental.py"

TEXT_SUFFIXES  = {".txt", ".md"}
PDF_SUFFIXES   = {".pdf"}
DOCX_SUFFIXES  = {".docx"}
PPTX_SUFFIXES  = {".pptx"}
EMAIL_SUFFIXES = {".eml"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
ALL_SUPPORTED  = TEXT_SUFFIXES | PDF_SUFFIXES | DOCX_SUFFIXES | PPTX_SUFFIXES | EMAIL_SUFFIXES | IMAGE_SUFFIXES

MEDIA_TYPES = {
    ".png":  "image/png",
    ".jpg":  "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif":  "image/gif",
    ".webp": "image/webp",
}

CLASSIFY_PROMPT = """\
You are classifying a document for a personal knowledge base.

Extract the following fields and respond with ONLY valid JSON (no markdown, no commentary):

{{
  "date": "YYYY-MM-DD or empty string if unknown",
  "category": "one of: NTA, DCC, DFB, ADAPT, Diotima, Paradigm, TBS, personal, other",
  "topic": "short topic label (5-10 words)",
  "title": "document title or descriptive title if none",
  "summary": "2-4 sentence summary of the document",
  "key_people": ["FirstName LastName"],
  "tags": ["tag1", "tag2"]
}}

Rules:
- category: NTA = transport/bus/rail/PSV/taxi, DCC = Dublin City Council,
  DFB = Dublin Fire Brigade, ADAPT = research/AI/academic,
  Diotima = strategy/consulting/Siobhan/Jonathan/Masa,
  Paradigm = technology, TBS = The Buildings Studio, personal = personal matters
- date: use the document's own date if present, otherwise leave empty
- key_people: only real named people, not generic roles; include email sender/recipients
- tags: 2-5 relevant keywords

Document content:
---
{content}
---"""


# ── LLM helpers ───────────────────────────────────────────────────────────────

def llm_call(messages: list, max_tokens: int = 512) -> str:
    resp = requests.post(
        LITELLM_URL,
        json={"model": MODEL, "messages": messages, "max_tokens": max_tokens},
        timeout=60,
    )
    resp.raise_for_status()
    return resp.json()["choices"][0]["message"]["content"]


def parse_json_response(raw: str) -> dict:
    raw = raw.strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)
    m = re.search(r"\{[\s\S]+\}", raw)
    if m:
        raw = m.group(0)
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        return {}


def apply_defaults(result: dict, today: str) -> dict:
    result.setdefault("date", today)
    result.setdefault("category", "other")
    result.setdefault("topic", "Document")
    result.setdefault("title", "Untitled")
    result.setdefault("summary", "")
    result.setdefault("key_people", [])
    result.setdefault("tags", [])
    if not result.get("date"):
        result["date"] = today
    return result


def classify_text(text: str, today: str) -> dict:
    prompt = CLASSIFY_PROMPT.format(content=text[:8000])
    raw = llm_call([{"role": "user", "content": prompt}])
    result = parse_json_response(raw)
    return apply_defaults(result, today)


# ── text extraction ───────────────────────────────────────────────────────────

def extract_text_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    parts = [page.extract_text() for page in reader.pages if page.extract_text()]
    return "\n\n".join(parts)


def extract_text_pdf_path(path: Path) -> str:
    return extract_text_pdf(path.read_bytes())


def extract_text_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_docx_path(path: Path) -> str:
    return extract_text_docx(path.read_bytes())


def extract_text_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(
                    para.text for para in shape.text_frame.paragraphs if para.text.strip()
                )
                if text:
                    parts.append(text)
        if parts:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(parts))
    return "\n\n".join(slides)


def extract_text_plain(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def strip_html(h: str) -> str:
    """Very basic HTML → plain text: strip tags, decode entities, collapse whitespace."""
    h = re.sub(r"<br\s*/?>", "\n", h, flags=re.IGNORECASE)
    h = re.sub(r"<p[^>]*>", "\n", h, flags=re.IGNORECASE)
    h = re.sub(r"<[^>]+>", "", h)
    h = html.unescape(h)
    h = re.sub(r"\n{3,}", "\n\n", h)
    return h.strip()


# ── email parsing ─────────────────────────────────────────────────────────────

def parse_email_date(date_str: str) -> str:
    """Parse email date header → YYYY-MM-DD, or empty on failure."""
    if not date_str:
        return ""
    from email.utils import parsedate_to_datetime
    try:
        return parsedate_to_datetime(date_str).strftime("%Y-%m-%d")
    except Exception:
        # Try simple regex
        m = re.search(r"\d{4}-\d{2}-\d{2}", date_str)
        return m.group(0) if m else ""


def name_from_address(addr: str) -> str:
    """'Rachel Buckley <buckler4@tcd.ie>' → 'Rachel Buckley'"""
    if not addr:
        return ""
    m = re.match(r"^([^<]+)<", addr.strip())
    if m:
        return m.group(1).strip()
    return addr.strip()


def extract_eml(path: Path) -> dict:
    """
    Parse a .eml file. Returns:
      {
        "subject": str,
        "from": str,
        "to": str,
        "date": str (YYYY-MM-DD),
        "body": str,          # plain text body
        "attachments": [{"name": str, "text": str}, ...]
      }
    """
    with open(path, "rb") as f:
        msg = email.message_from_binary_file(f, policy=email_policy.default)

    subject   = msg.get("subject", "")
    from_addr = msg.get("from", "")
    to_addr   = msg.get("to", "")
    date_str  = parse_email_date(msg.get("date", ""))

    body_plain = ""
    body_html  = ""
    attachments = []

    # Inline image content-types to skip (decorative Outlook images etc.)
    skip_types = {"image/png", "image/jpeg", "image/gif", "image/webp"}

    for part in msg.walk():
        ct   = part.get_content_type()
        disp = part.get_content_disposition() or ""
        name = part.get_filename() or ""

        if ct == "text/plain" and disp != "attachment" and not body_plain:
            try:
                body_plain = part.get_content()
            except Exception:
                body_plain = part.get_payload(decode=True).decode("utf-8", errors="replace")

        elif ct == "text/html" and disp != "attachment" and not body_html:
            try:
                body_html = part.get_content()
            except Exception:
                body_html = part.get_payload(decode=True).decode("utf-8", errors="replace")

        elif disp == "attachment" and name:
            data = part.get_payload(decode=True)
            if not data:
                continue
            suffix = Path(name).suffix.lower()
            att_text = ""
            if suffix == ".pdf":
                try:
                    att_text = extract_text_pdf(data)
                except Exception as e:
                    att_text = f"[PDF extraction failed: {e}]"
            elif suffix in {".docx"}:
                try:
                    att_text = extract_text_docx(data)
                except Exception as e:
                    att_text = f"[DOCX extraction failed: {e}]"
            elif suffix in {".pptx"}:
                try:
                    import tempfile
                    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                        tmp.write(data)
                        tmp_path = Path(tmp.name)
                    att_text = extract_text_pptx(tmp_path)
                    tmp_path.unlink()
                except Exception as e:
                    att_text = f"[PPTX extraction failed: {e}]"
            elif suffix in {".txt", ".md"}:
                att_text = data.decode("utf-8", errors="replace")
            else:
                att_text = f"[Attachment type {suffix} not extracted]"

            attachments.append({"name": name, "text": att_text})

        elif ct in skip_types and disp == "inline":
            pass  # Skip decorative inline images

    body = body_plain if body_plain.strip() else strip_html(body_html)

    return {
        "subject":     subject,
        "from":        from_addr,
        "to":          to_addr,
        "date":        date_str,
        "body":        body,
        "attachments": attachments,
    }


def process_eml(path: Path, today: str) -> tuple:
    """
    Parse .eml and produce (meta, body_markdown, extra_frontmatter).
    body_markdown is formatted for the KB doc.
    """
    parsed = extract_eml(path)

    # Build text for classification: header context + body + attachment text
    classify_parts = [
        f"Subject: {parsed['subject']}",
        f"From: {parsed['from']}",
        f"To: {parsed['to']}",
        f"Date: {parsed['date']}",
        "",
        parsed["body"],
    ]
    for att in parsed["attachments"]:
        if att["text"] and not att["text"].startswith("["):
            classify_parts += ["", f"--- Attachment: {att['name']} ---", att["text"][:3000]]

    classify_text_str = "\n".join(classify_parts)

    meta = classify_text(classify_text_str, today)

    # Use email date if available (more reliable than LLM extraction)
    if parsed["date"]:
        meta["date"] = parsed["date"]

    # Use subject as title fallback
    if meta.get("title") in ("Untitled", "") and parsed["subject"]:
        meta["title"] = parsed["subject"]

    # Merge sender/recipient into key_people
    from_name = name_from_address(parsed["from"])
    to_names  = [name_from_address(a.strip()) for a in parsed["to"].split(",")]
    extra_people = {n for n in [from_name] + to_names
                    if n and "eoin" not in n.lower() and "noval" not in n.lower()}
    existing = {p.lower() for p in meta.get("key_people", [])}
    for p in sorted(extra_people):
        if p.lower() not in existing:
            meta["key_people"].append(p)

    # Build the formatted body for the KB doc
    body_parts = [
        f"**From:** {parsed['from']}",
        f"**To:** {parsed['to']}",
        f"**Date:** {parsed['date'] or 'unknown'}",
        f"**Subject:** {parsed['subject']}",
        "",
        parsed["body"].strip(),
    ]

    for att in parsed["attachments"]:
        body_parts += ["", f"---", f"### Attachment: {att['name']}", ""]
        if att["text"]:
            body_parts.append(att["text"])
        else:
            body_parts.append("*[No text extracted]*")

    body_md = "\n".join(body_parts)

    extra_fm = {
        "type":       "email",
        "email_from": parsed["from"],
        "email_to":   parsed["to"],
    }

    return meta, body_md, extra_fm


# ── image classification ───────────────────────────────────────────────────────

def classify_image(path: Path, today: str) -> tuple:
    media_type = MEDIA_TYPES.get(path.suffix.lower(), "image/jpeg")
    image_data = base64.standard_b64encode(path.read_bytes()).decode("utf-8")

    messages = [{
        "role": "user",
        "content": [
            {
                "type": "image_url",
                "image_url": {"url": f"data:{media_type};base64,{image_data}"},
            },
            {
                "type": "text",
                "text": (
                    "First describe what this image contains in detail — extract all visible text verbatim.\n\n"
                    "Then on a new line output ONLY valid JSON:\n"
                    '{"date":"YYYY-MM-DD or empty","category":"NTA/DCC/DFB/ADAPT/Diotima/Paradigm/TBS/personal/other",'
                    '"topic":"5-10 word topic","title":"title","summary":"2-4 sentence summary",'
                    '"key_people":[],"tags":[]}'
                ),
            },
        ],
    }]

    raw = llm_call(messages, max_tokens=1024)
    json_match = re.search(r"\{[\s\S]+\}", raw)
    description = raw[:json_match.start()].strip() if json_match else raw
    result = parse_json_response(json_match.group(0) if json_match else "{}")
    result.setdefault("summary", description[:600])
    result.setdefault("title", path.stem)

    return description, apply_defaults(result, today)


# ── KB output ─────────────────────────────────────────────────────────────────

def slugify(s: str) -> str:
    return re.sub(r"[^a-z0-9-]", "", s.lower().replace(" ", "-").replace("'", ""))[:60]


def write_kb_markdown(meta: dict, body: str, source_filename: str,
                      extra_frontmatter: dict = None) -> Path:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)

    date_str = meta["date"]
    category = meta["category"]
    slug     = slugify(meta["topic"])
    filename = f"{date_str}_{category}_{slug}.md"
    out_path = DOCS_DIR / filename

    if out_path.exists():
        out_path = DOCS_DIR / f"{out_path.stem}_{datetime.now().strftime('%H%M%S')}{out_path.suffix}"

    people_yaml = "\n".join(f'  - "{p}"' for p in meta["key_people"])
    tags_yaml   = "\n".join(f'  - "{t}"' for t in meta["tags"])

    extra_lines = ""
    if extra_frontmatter:
        for k, v in extra_frontmatter.items():
            safe_v = str(v).replace('"', "'")
            extra_lines += f'\n{k}: "{safe_v}"'

    content = f"""---
title: "{meta['title'].replace('"', "'")}"
date: {date_str}
category: {category}
topic: "{meta['topic'].replace('"', "'")}"
people:
{people_yaml if people_yaml else "  []"}
tags:
{tags_yaml if tags_yaml else "  []"}
source_file: "{source_filename}"{extra_lines}
---

## Summary

{meta['summary']}

## Content

{body}
"""

    out_path.write_text(content, encoding="utf-8")
    return out_path


# ── upload trigger ────────────────────────────────────────────────────────────

def trigger_upload():
    if not UPLOAD_SCRIPT.exists():
        return
    try:
        subprocess.Popen(
            ["/usr/local/bin/python3", str(UPLOAD_SCRIPT)],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        print("  Upload triggered (running in background)")
    except Exception as e:
        print(f"  Upload trigger failed (non-fatal): {e}")


# ── main file processor ───────────────────────────────────────────────────────

def process_file(path: Path, today: str) -> bool:
    suffix = path.suffix.lower()
    print(f"  Processing: {path.name}")

    try:
        extra_fm = None

        if suffix in TEXT_SUFFIXES:
            text = extract_text_plain(path)
            meta = classify_text(text, today)
            body = text

        elif suffix in PDF_SUFFIXES:
            text = extract_text_pdf_path(path)
            if not text.strip():
                print("    WARNING: No text extracted from PDF (may be scanned image)")
                text = f"[PDF file: {path.name} — no extractable text]"
            meta = classify_text(text, today)
            body = text

        elif suffix in DOCX_SUFFIXES:
            text = extract_text_docx_path(path)
            meta = classify_text(text, today)
            body = text

        elif suffix in PPTX_SUFFIXES:
            text = extract_text_pptx(path)
            if not text.strip():
                text = f"[Presentation: {path.name} — no extractable text]"
            meta = classify_text(text, today)
            body = text

        elif suffix in EMAIL_SUFFIXES:
            meta, body, extra_fm = process_eml(path, today)

        elif suffix in IMAGE_SUFFIXES:
            description, meta = classify_image(path, today)
            body = description

        else:
            print(f"    SKIP: unsupported type {suffix}")
            return False

        out = write_kb_markdown(meta, body, path.name, extra_fm)
        print(f"    → {out.relative_to(Path.home())}")
        print(f"    Category: {meta['category']} | Topic: {meta['topic']}")

        DONE_DIR.mkdir(parents=True, exist_ok=True)
        dest = DONE_DIR / path.name
        if dest.exists():
            dest = DONE_DIR / f"{path.stem}_{datetime.now().strftime('%H%M%S')}{path.suffix}"
        shutil.move(str(path), str(dest))

        return True

    except Exception as e:
        import traceback
        print(f"    ERROR: {e}")
        traceback.print_exc()
        return False


def main():
    INBOX_DIR.mkdir(parents=True, exist_ok=True)

    files = [
        f for f in INBOX_DIR.iterdir()
        if f.is_file() and f.suffix.lower() in ALL_SUPPORTED
    ]

    if not files:
        print("Inbox is empty — nothing to process.")
        return

    print(f"Found {len(files)} file(s) to process...")
    today = datetime.now().strftime("%Y-%m-%d")

    processed = 0
    for f in sorted(files):
        if process_file(f, today):
            processed += 1

    if processed:
        print(f"\nProcessed {processed}/{len(files)} files → triggering upload...")
        trigger_upload()

    print("Done.")


if __name__ == "__main__":
    main()
